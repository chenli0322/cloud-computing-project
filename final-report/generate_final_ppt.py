"""
Generates the Final Presentation .pptx for the NYU Cloud Computing project.

Two parts back-to-back:
  Part 1 (~20 min, 8 slides): ArchNav legacy migration to AWS EC2 via Docker
  Part 2 (~25 min, 13 slides): Decentralized Health Monitoring System

Run from final-report/:
    python generate_final_ppt.py

Output: final-report/Cloud_Computing_Final_Presentation.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


HERE = Path(__file__).parent
OUT = HERE / "Cloud_Computing_Final_Presentation.pptx"

# Brand palette (NYU-ish + dark accent)
NYU_VIOLET = RGBColor(0x57, 0x0E, 0xA0)
DARK_BG = RGBColor(0x0F, 0x14, 0x1F)
LIGHT_TEXT = RGBColor(0xE6, 0xEA, 0xF2)
MUTED = RGBColor(0x8B, 0x95, 0xAB)
ACCENT = RGBColor(0x6A, 0xA6, 0xFF)
GREEN = RGBColor(0x4E, 0xC9, 0x6F)
RED = RGBColor(0xE8, 0x5C, 0x57)


def add_dark_bg(slide, prs):
    """Fill the slide with a dark background rectangle (behind everything)."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_title(slide, text, top=Inches(0.4), size=32, color=LIGHT_TEXT, bold=True):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_subtitle(slide, text, top=Inches(1.1), size=18, color=MUTED):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = True
    return box


def add_bullets(slide, bullets, top=Inches(1.7), left=Inches(0.7),
                width=Inches(12), height=Inches(5.5), size=18,
                color=LIGHT_TEXT, accent_first_word=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        if isinstance(b, tuple):
            # (level, text) or (level, text, color)
            if len(b) == 2:
                level, text = b
                col = color
            else:
                level, text, col = b
            p.level = level
        else:
            text = b
            col = color
        r = p.add_run()
        r.text = "•  " + text if (isinstance(b, str) or (isinstance(b, tuple) and b[0] == 0)) else "–  " + text
        r.font.size = Pt(size if (isinstance(b, str) or (isinstance(b, tuple) and b[0] == 0)) else size - 2)
        r.font.color.rgb = col
    return box


def add_footer(slide, prs, text):
    box = slide.shapes.add_textbox(
        Inches(0.6), prs.slide_height - Inches(0.45),
        Inches(12), Inches(0.3),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_accent_bar(slide, top=Inches(1.05), color=ACCENT, width_in=2.5):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(width_in), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_table(slide, data, left, top, width, height,
              header_color=NYU_VIOLET, body_color=DARK_BG,
              header_text=LIGHT_TEXT, body_text=LIGHT_TEXT, font_size=14):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else body_color
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(font_size)
            run.font.color.rgb = header_text if r == 0 else body_text
            run.font.bold = (r == 0)
    return tbl


def add_speaker_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ----------------------------------------------------------------------
# Build the deck
# ----------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


def new_slide(title=None, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)
    add_dark_bg(slide, prs)
    if title is not None:
        add_title(slide, title)
    if subtitle is not None:
        add_subtitle(slide, subtitle)
        add_accent_bar(slide, top=Inches(1.55))
    else:
        add_accent_bar(slide, top=Inches(1.05))
    return slide


# ============ TITLE SLIDE ============
s = prs.slides.add_slide(blank_layout)
add_dark_bg(s, prs)
add_title(
    s,
    "Cloud Computing — Spring 2026 Final Project",
    top=Inches(2.4), size=40,
)
add_subtitle(
    s,
    "Part 1: ArchNav Legacy Migration via Docker on AWS EC2",
    top=Inches(3.4), size=20, color=ACCENT,
)
add_subtitle(
    s,
    "Part 2: Decentralized Health Monitoring System (P2P + IoT + ML + Blockchain)",
    top=Inches(3.85), size=20, color=ACCENT,
)
# Author block
box = s.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12), Inches(1.5))
tf = box.text_frame
for line, sz, col in [
    ("Chen Li (NetID: cl5725) — solo team", 18, LIGHT_TEXT),
    ("Special CS Topic — Cloud Computing, Section 026", 16, MUTED),
    ("Prof. Jean-Claude Franchitti — NYU Tandon", 16, MUTED),
    ("April 30, 2026", 14, MUTED),
]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    r = p.add_run()
    r.text = line
    r.font.size = Pt(sz)
    r.font.color.rgb = col
add_speaker_notes(s, """Welcome. I'm Chen Li, presenting solo for both Part 1 and Part 2 of the cloud computing project. The plan: roughly 20 minutes on Part 1 (the ArchNav legacy migration) and 25 minutes on Part 2 (the decentralized health monitoring system). At the end I'll demo the live system.""")

# ============================================================
# ============ PART 1 — ArchNav Migration ====================
# ============================================================

# --- Part 1 cover ---
s = new_slide("Part 1 — Legacy Cloud Migration", "ArchNav (Oracle ADF + Fortress) → Docker → AWS EC2")
add_bullets(s, [
    "Cloud-based migration: zero source-code changes",
    "5-component Java legacy stack containerised into 4 Docker services",
    "Deployed on AWS EC2 (us-east-1, Ubuntu 24.04), accessible at public IP",
    "All services running on a single bridge network (`archnav-net`)",
    "One-command rebuild via deploy.sh (~30 min from cold metal to live demo)",
], top=Inches(1.95))
add_footer(s, prs, "Part 1 — slide 1 / 8")
add_speaker_notes(s, """Part 1 is a classic cloud-based migration: the goal is to take an existing Java application that nobody is willing to rewrite, and put it in the cloud without touching the source. We containerise the dependencies, deploy with Docker Compose, and run on AWS EC2. The whole thing rebuilds in about 30 minutes.""")

# --- Slide P1-2: What is ArchNav ---
s = new_slide("ArchNav: What Are We Migrating?")
add_bullets(s, [
    "ArchNav = Oracle ADF web application for architecture-knowledge navigation",
    "Built circa 2017 against very specific, very fragile dependencies:",
    (1, "JDK 1.7.0_80 (exact patch version)"),
    (1, "Oracle JDeveloper 11.1.2.4 + ADF Essentials 11.1.2.4"),
    (1, "Glassfish 3.1.2 (HTTP port 9999, admin 4848)"),
    (1, "MySQL 5.7 (schema: archemy)"),
    (1, "Apache Fortress 2.0.0-RC1 + ApacheDS 2.0.0-M23 + Tomcat 8.5.11"),
    "Two architectural stacks: business (Glassfish/ADF/MySQL) + security (Fortress/LDAP/Tomcat)",
    "RBAC enforced via Fortress ELResolver inside the Glassfish app",
], top=Inches(1.35), size=16)
add_footer(s, prs, "Part 1 — slide 2 / 8")
add_speaker_notes(s, """The application is small but its dependency surface is brittle: JDK 7 specifically, JDeveloper 11.1.2.4 specifically, Glassfish 3.1.2 specifically. Real-world legacy apps often look like this — they cannot move forward in version because nobody has budget to test the upgrade. That's exactly the situation cloud-based migration is designed for.""")

# --- Slide P1-3: Two-stack architecture (current state) ---
s = new_slide("Original Architecture: Two-Stack Windows Install")
# Two boxes
arch_box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
tf = arch_box.text_frame
tf.word_wrap = True
text = """SECURITY STACK (Apache Fortress + ApacheDS + Tomcat)
   Fortress Web UI       :8080/fortress-web   (test/password)
   Fortress REST API     :8080/fortress-rest
   Fortress Core         (Maven-built)
   Apache Tomcat 8.5.11  :8080
   ApacheDS LDAP         :389  (uid=admin,ou=system / secret)

BUSINESS STACK (Glassfish + ADF + MySQL)
   ArchNav Web UI        :9999/archemy/faces/login.jspx
   Oracle ADF 11.1.2.4   (calls Fortress via ELResolver)
   JDeveloper 11.1.2.4   (build+deploy to EAR)
   Glassfish 3.1.2       :9999  admin :4848  (Claude1960%)
   MySQL 5.7             :3306  schema=archemy  (archemy/archemydb1960%)

   Stack 1 enforces authn/authz; Stack 2 hosts the business logic.
   Cross-stack call: ADF code → ELResolver → Fortress REST → ApacheDS."""
for i, line in enumerate(text.split("\n")):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(13)
    r.font.color.rgb = LIGHT_TEXT if not line.startswith("   ") else MUTED
add_footer(s, prs, "Part 1 — slide 3 / 8")
add_speaker_notes(s, """Two stacks. Security on the left — Apache Fortress for RBAC, ApacheDS as the LDAP backend, Tomcat hosting the Fortress REST and Web. Business on the right — Glassfish 3.1.2 hosting the ADF application, MySQL 5.7 holding the data. The two stacks talk through Fortress ELResolver calls embedded in the ADF web pages.""")

# --- Slide P1-4: Future state (containerised) ---
s = new_slide("Target Architecture: 4 Containers + 1 Bridge Network")
add_bullets(s, [
    "Dockerfile.mysql       — MySQL 5.7 + archemy schema (init.sql)",
    "Dockerfile.apacheds    — ApacheDS 2.0.0-M23, Fortress LDAP schema preloaded",
    "Dockerfile.tomcat      — Tomcat 8.5.11 + fortress-rest.war + fortress-web.war",
    "Dockerfile.glassfish   — JDK 7 + Glassfish 3.1.2 + ADF Essentials + ArchNav EAR",
    "docker-compose.yml     — 4 services on bridge network `archnav-net`",
    "  • Inter-service hostnames replace localhost (apacheds, mysql, tomcat, glassfish)",
    "  • External ports: 3306, 8080, 9999, 389, 4848",
    "Single-command bring-up: `docker compose up --build -d`  (~6 min cold build on EC2)",
], top=Inches(1.35), size=15)
add_footer(s, prs, "Part 1 — slide 4 / 8")
add_speaker_notes(s, """Each of the 5 components becomes a container. The 4 containers communicate over a Docker bridge network using service names as hostnames — so the Fortress configs that used to point at "localhost" now point at "apacheds", "mysql", etc. That hostname change is one of the most common pitfalls in cloud-based migration.""")

# --- Slide P1-5: AWS EC2 deployment ---
s = new_slide("AWS EC2 Deployment")
add_bullets(s, [
    "Instance type: t2.small (1 vCPU, 2 GB RAM) — sufficient for the demo",
    "AMI: Ubuntu Server 24.04 LTS, region us-east-1, 20 GiB gp3 storage",
    "Security group: SSH(22) + 8080 + 9999 + 4848 inbound",
    "Provisioning: SCP the docker/ tree (~278 MB) → `apt install docker.io` → `docker compose up`",
    "First-build time: ~6 min on a fresh instance",
    "Live URL during evaluation: http://<public-ip>:9999/archemy/faces/login.jspx",
    "Reproducible bring-up via deploy.sh (one-command); rebuild script committed in repo",
], top=Inches(1.35), size=16)
add_footer(s, prs, "Part 1 — slide 5 / 8")
add_speaker_notes(s, """The cloud target is AWS EC2. t2.small was actually a tight fit but it works for the demo. The whole bring-up — SSH to the box, install Docker, SCP the docker directory, run docker compose up — is wrapped in a deploy.sh script so the evaluation rebuild on April 29 is a single command.""")

# --- Slide P1-6: Engineering challenges ---
s = new_slide("Engineering Challenges Encountered")
add_bullets(s, [
    "JDK 7 RPM packaging: extracted .jar files were in `.pack` format → added `unpack200` step in Dockerfile.glassfish",
    "Glassfish self-signed certs expired in 2022 → removed `enable-secure-admin`, used `--secure=false`",
    "Fortress WAR filename mismatch: Maven produces `fortress-rest.war`, Dockerfile copied `fortress-rest-2.0.0-RC1.war` → fixed COPY paths",
    "Fortress LDAP config scattered in 4 nested locations (lib JAR, expanded WAR JAR, __internal EAR, /opt EAR) → all four had to be patched with `host=apacheds, dc=example,dc=com, uid=admin,ou=system`",
    "Docker-Desktop WSL2 crash + C: drive full → WSL distro moved to F: drive",
    "All issues documented in install_log.md with exact commands for reproducibility",
], top=Inches(1.35), size=14)
add_footer(s, prs, "Part 1 — slide 6 / 8")
add_speaker_notes(s, """The interesting bit of Part 1 was the debug story. JDK 7 RPMs ship .pack files instead of .jars; Glassfish ships SSL certs that expired in 2022; the Fortress LDAP config exists in four different nested archives that all need to be patched. Each of these is the kind of thing that takes one engineer one afternoon and is invisible in the final running system.""")

# --- Slide P1-7: Result — before/after ---
s = new_slide("Migration Result: Cloud-Based, No Code Changes")
add_bullets(s, [
    "ZERO source code changes to ArchNav, ADF, Fortress, or MySQL — pure environmental migration",
    "BEFORE: Single Windows host with 17-step manual install (~8 hours skilled engineer time)",
    "AFTER: `git clone && cd docker && docker compose up --build -d`  (~6 min)",
    "Disposable: `docker compose down -v` returns the host to clean state",
    "Reproducible: same Dockerfiles produce same containers on developer laptop and EC2 instance",
    "Ready for further evolution: cloud-enabled refactor (Part 2-style cloud APIs) becomes feasible",
], top=Inches(1.35), size=15)
add_footer(s, prs, "Part 1 — slide 7 / 8")
add_speaker_notes(s, """The deliverable is exactly the goal of cloud-based migration: same code, different deployment surface. The benefit isn't performance or feature — it's that an 8-hour 17-step Windows install becomes a 6-minute one-line command, on disposable hardware, reproducibly.""")

# --- Slide P1-8: Live demo ---
s = new_slide("Part 1 Live Demo")
add_bullets(s, [
    "Open public IP in browser: http://<ec2-ip>:9999/archemy/faces/login.jspx",
    "Login as Admin user (tom@oracle.com / oracle123) → ArchNav homepage",
    "Login as NormalUser (john@oracle.com / oracle123) → restricted view (Fortress RBAC enforces it)",
    "Show Fortress Web at http://<ec2-ip>:8080/fortress-web (test / password) — see the role/permission model",
    "SSH demo: `docker compose ps`, `docker compose logs glassfish | tail`",
    "→ Transition to Part 2",
], top=Inches(1.55), size=17)
add_footer(s, prs, "Part 1 — slide 8 / 8")
add_speaker_notes(s, """For the demo I'll log in as both an Admin and a NormalUser to show that Fortress RBAC is actually enforcing distinct views. I'll also show the Fortress Web admin console where the role/permission model is visible. Then I'll switch over to Part 2.""")


# ============================================================
# ============ PART 2 — Health Monitor =======================
# ============================================================

# --- Part 2 cover ---
s = new_slide("Part 2 — Decentralized Health Monitoring System",
               "Hybrid Web2 + Web3 — IoT, ML, Blockchain, broker-less P2P")
add_bullets(s, [
    "5 cooperating processes, broker-less P2P mesh — no central message router",
    "3 cloud platforms actively engaged: Azure IoT Hub + AWS S3 + Ethereum Sepolia",
    "ML model artefact lives in S3, downloaded by ML node at startup",
    "Every flagged anomaly: archived to S3, hash anchored on-chain, dashboard shows both links",
    "End-to-end auditable: pull S3 → recompute SHA-256 → matches on-chain hash",
], top=Inches(1.95))
add_footer(s, prs, "Part 2 — slide 1 / 13")
add_speaker_notes(s, """Part 2 is the innovative system. The headline is: a fully decentralised health-monitoring pipeline running across three real cloud platforms, broker-less, with cryptographic auditability between the off-chain raw record and the on-chain hash. I'll walk through the architecture, show the live numbers from yesterday's rehearsal, and finish with the concrete audit-chain verification.""")

# --- Slide P2-2: Problem & goals ---
s = new_slide("Problem & Goals")
add_bullets(s, [
    "Today's IoT health platforms are single-tenant: one operator holds the data and the audit log",
    "Two failure modes: (1) the operator can silently rewrite history, (2) the operator going down takes everyone down",
    "Goal: distribute processing across independent peers + anchor the audit log on a public chain",
    (1, "G1 — broker-less P2P with bootstrap discovery and master election"),
    (1, "G2 — IoT simulator with ground-truth labels for ML evaluation"),
    (1, "G3 — unsupervised anomaly detector with reasonable recall on synthetic data"),
    (1, "G4 — deploy a real smart contract on a real public chain (not a local simulator)"),
    (1, "G5 — single dashboard fusing sensor stream, ML verdicts, peer state, on-chain proofs"),
], top=Inches(1.35), size=15)
add_footer(s, prs, "Part 2 — slide 2 / 13")
add_speaker_notes(s, """The motivation: today's healthcare IoT platforms are single-tenant and single-trust. We want auditability — a verifier should be able to prove a record hasn't been tampered with without trusting the operator. That's exactly what a public chain anchor gives you. The five goals on the slide are what the rest of this part will demonstrate.""")

# --- Slide P2-3: Architecture diagram (text) ---
s = new_slide("Architecture: 5 Roles + 3 Clouds")
arch_box = s.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.5))
tf = arch_box.text_frame
tf.word_wrap = True
arch = """                                ┌──────────────┐
                                │ Azure IoT Hub │  ← MQTT 8883/TLS, F1 free tier
                                └──────▲───────┘
                                       │ telemetry
   ┌────────────┐  MSG_SENSOR  ┌───────┴──────┐                     ┌──────────┐
   │ IoT (sim)  │─────────────▶│   ML node    │◀──── load model ───│  AWS S3  │
   └────────────┘              └──────┬───────┘                     └──────────┘
                                      │ MSG_ANOMALY                       ▲
                                      ▼                                   │
                              ┌────────────────┐  PUT raw JSON ───────────┘
                              │   BC node      │                          │
                              │ NonceMgr+S3    │  logAnomaly(hash, …)     │
                              └──────┬─────────┘                          │
                                     ▼                                    │
                            ┌────────────────┐                            │
                            │  Sepolia chain │  block confirmation        │
                            └──────┬─────────┘                            │
                                   │ MSG_BC_LOGGED (tx + s3_uri)          │
                                   ▼                                      │
                          ┌──────────────────┐                            │
                          │     Dashboard    │ ─ aiohttp WS ─ Browser     │
                          │  (Etherscan + S3 │      ▲                     │
                          │   raw-JSON link) │      │                     │
                          └──────────────────┘      └─ S3 raw JSON ───────┘

   Bootstrap (5th role) is off the data path: discovery only.
   Cloud services in use: Azure IoT Hub + AWS S3 + Ethereum Sepolia (3 PaaS / 3 clouds)."""
for i, line in enumerate(arch.split("\n")):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(11)
    r.font.color.rgb = LIGHT_TEXT
add_footer(s, prs, "Part 2 — slide 3 / 13")
add_speaker_notes(s, """The architecture has five roles communicating over a broker-less P2P mesh. The IoT node mirrors every reading to Azure IoT Hub. The ML node downloads its model from AWS S3. The BC node, before submitting to Sepolia, archives the raw anomaly JSON to S3. The dashboard fuses everything for the browser. Three clouds, three roles for them.""")

# --- Slide P2-4: P2P framework decision ---
s = new_slide("P2P Layer: BTP / libp2p Decision (Disclosure)")
add_bullets(s, [
    "Midterm proposal claimed BTP framework. After investigation:",
    (1, "BTP public Python/Java repos last updated 2018, archived"),
    (1, "py-libp2p alpha quality, incomplete Kademlia DHT"),
    (1, "Migrating ~950 LOC four days before demo = unacceptable schedule risk"),
    "Built equivalent broker-less overlay (~280 lines, p2p-network/peer_node.py)",
    "All architectural properties present:",
    (1, "Broker-less direct-connection mesh ✓"),
    (1, "Bootstrap-based peer discovery ✓ (any peer answers MSG_HELLO, not just bootstrap)"),
    (1, "Bully leader election ✓ — _start_election(), bully algorithm"),
    (1, "Liveness keepalive ✓ — MSG_PING/MSG_PONG every 15s"),
    (1, "Wire protocol: 4-byte big-endian length prefix + UTF-8 JSON"),
    "Disclosed in §3.6 of final report — substitution driven by upstream readiness, not scope reduction",
], top=Inches(1.35), size=14)
add_footer(s, prs, "Part 2 — slide 4 / 13")
add_speaker_notes(s, """Honest disclosure: I said BTP at the midterm. After investigation, both BTP and py-libp2p were not production-ready in the project timeline. So I built an equivalent broker-less overlay with the same properties. This is in section 3.6 of the report. The substitution is for upstream-library reasons, not because we reduced scope.""")

# --- Slide P2-5: IoT + Azure IoT Hub ---
s = new_slide("IoT Layer + Azure IoT Hub")
add_bullets(s, [
    "Synthetic sensor: heart rate, body temp, SpO₂ — clipped Gaussians around resting-adult ranges",
    "Injected anomaly rate: 8% by default — four kinds (hr_high, hr_low, temp_high, spo2_low)",
    "Each reading: (1) broadcast as MSG_SENSOR over P2P, (2) mirrored to Azure IoT Hub via MQTT 8883/TLS",
    "Azure auth: SAS token over `<host>/devices/<device-id>` using HMAC-SHA256 (no extra SDK)",
    "Free tier F1: 8 000 messages/day, $0/month — comfortably handles demo traffic",
    "Connection-state logging: 'azure: published N messages so far' every 10 reads",
    "Last rehearsal (2026-04-26): 310 messages published in ~8 minutes, 100% success rate",
], top=Inches(1.35), size=15)
add_footer(s, prs, "Part 2 — slide 5 / 13")
add_speaker_notes(s, """The IoT layer publishes to two destinations: the internal P2P mesh, and Azure IoT Hub. The Azure path satisfies the IoT-PaaS requirement. We use the free tier F1 — 8000 messages a day, no charge — and during yesterday's rehearsal we published 310 messages with no failures.""")

# --- Slide P2-6: ML + S3 model storage ---
s = new_slide("ML Layer + S3 Model Storage")
add_bullets(s, [
    "Algorithm: Isolation Forest (sklearn, n_estimators=200, contamination=0.1)",
    "Training set: 5 000 synthetic healthy samples; held-out test: 1 500 (seed=99)",
    "Reported metrics on test set:",
], top=Inches(1.35), size=15)
add_table(s, [
    ["Class", "Precision", "Recall", "F1", "Support"],
    ["Normal", "0.99", "0.99", "0.99", "1 350"],
    ["Anomaly", "0.90", "0.98", "0.94", "150"],
    ["Overall accuracy", "—", "—", "0.99", "1 500"],
], left=Inches(0.7), top=Inches(2.95), width=Inches(8.5), height=Inches(1.6))
add_bullets(s, [
    "Model artefact (1.6 MB joblib) uploaded to s3://chenli-cloud-final-2026/models/anomaly_model.joblib",
    "ml_node.py reads MODEL_S3_URI from .env, downloads at startup via boto3, joblib.load from temp file",
    "Verified at rehearsal: ML node logged 'Found credentials in environment variables' + S3 download",
], top=Inches(4.85), size=14)
add_footer(s, prs, "Part 2 — slide 6 / 13")
add_speaker_notes(s, """ML is an Isolation Forest — unsupervised, fast, no tuning. On the held-out test set we get 99% accuracy with 98% anomaly recall. The trained joblib artefact lives in S3, and the ML node downloads it at startup via boto3. This makes the model "on cloud" — training can happen anywhere, but the consumed artefact is always pulled from S3.""")

# --- Slide P2-7: Smart contract ---
s = new_slide("Smart Contract: HealthLog.sol on Sepolia")
code = """// SPDX-License-Identifier: MIT
pragma solidity ^0.8.20;

contract HealthLog {
    struct Entry {
        bytes32 eventHash;     // sha256 of the off-chain JSON event
        uint256 timestamp;
        address reporter;
        string  deviceId;
        string  anomalyKind;
    }
    Entry[] public entries;
    mapping(bytes32 => uint256) public indexOfHash;
    event AnomalyLogged(bytes32 indexed eventHash, address indexed reporter,
                        string deviceId, string anomalyKind, uint256 timestamp);

    function logAnomaly(bytes32 eventHash, string calldata deviceId,
                        string calldata anomalyKind) external returns (uint256 id) {
        require(indexOfHash[eventHash] == 0, "duplicate event");
        entries.push(Entry({...}));
        id = entries.length;
        indexOfHash[eventHash] = id;
        emit AnomalyLogged(eventHash, msg.sender, deviceId, anomalyKind, block.timestamp);
    }
}"""
box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8), Inches(4.5))
tf = box.text_frame
for i, line in enumerate(code.split("\n")):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(11)
    r.font.color.rgb = LIGHT_TEXT
add_bullets(s, [
    "Deployed: 0x89983910f6AE98Ea081356148B433cA3C6de283B",
    "Optimizer: 200 runs, Solidity 0.8.24",
    "Append-only log + duplicate-rejection",
    "Events indexed on hash + reporter for cheap querying",
    "View funcs: exists(hash), getEntry(id), total()",
], left=Inches(8.9), top=Inches(1.6), width=Inches(4.0), size=12)
add_footer(s, prs, "Part 2 — slide 7 / 13")
add_speaker_notes(s, """The contract is deliberately small. An append-only entries array, a hash → index mapping for duplicate rejection, and a single mutating function. The whole contract is 67 lines including comments. This is what gets deployed on Sepolia at the address shown.""")

# --- Slide P2-8: Sepolia migration ---
s = new_slide("Why Sepolia (Not Ganache as Originally Proposed)")
add_bullets(s, [
    "Midterm proposed Ganache — local in-memory Ethereum simulator",
    "Instructor's 2026-04-16 guidance: \"go for a real public chain so the project demonstrates real PaaS use\"",
    "Switched to Sepolia public testnet:",
    (1, "Real signing, real gas accounting, real block-confirmation latency"),
    (1, "Public verifiability — Etherscan resolves the contract and every transaction"),
    (1, "Forced confrontation with real-world frictions (nonce contention, RPC rate limits)"),
    "Free Sepolia ETH from sepoliafaucet.com; ~0.000167 ETH per logAnomaly transaction",
    "STRICT UPGRADE relative to midterm proposal — driven by instructor guidance, not scope reduction",
], top=Inches(1.35), size=15)
add_footer(s, prs, "Part 2 — slide 8 / 13")
add_speaker_notes(s, """At the midterm I proposed Ganache, which is a local Ethereum simulator. After Session 9 on April 16, the instructor said go for a real public chain. So I switched to Sepolia, which means real transactions, real gas, real block times, real Etherscan links. This is documented in §3.5 and §6.7 of the report. It's the second of two material differences from the midterm proposal — both upgrades.""")

# --- Slide P2-9: BC submitter (the engineering-heavy slide) ---
s = new_slide("BC Submitter: The Engineering-Heavy Component")
add_bullets(s, [
    "On every MSG_ANOMALY:",
    (1, "1️⃣  S3Archiver.archive() — PUT raw JSON to s3://<bucket>/anomalies/<hash>.json"),
    (1, "2️⃣  NonceManager.get_nonce() — asyncio.Lock + cached next_nonce, rollback on failure"),
    (1, "3️⃣  contract.functions.logAnomaly(hash, deviceId, kind).build_transaction(...)"),
    (1, "4️⃣  signed=acct.sign_transaction(...); send_raw_transaction; wait_for_receipt(180s)"),
    (1, "5️⃣  if rcpt.status != 1: raise RuntimeError → no MSG_BC_LOGGED"),
    (1, "6️⃣  broadcast MSG_BC_LOGGED with tx_hash + block + gas_used + s3_uri"),
    "EIP-1559 params: gas=250 000, maxFeePerGas=5 gwei, maxPriorityFeePerGas=1 gwei",
    "S3 archival happens BEFORE chain submission — off-chain object exists when on-chain hash becomes permanent",
], top=Inches(1.35), size=14)
add_footer(s, prs, "Part 2 — slide 9 / 13")
add_speaker_notes(s, """The BC node is the engineering-heavy bit because it sits at the Web2-to-Web3 boundary. The order matters: archive to S3 first, then anchor on chain. If the order were reversed and the chain commit succeeded but the S3 archive failed, you'd have an immutable hash with no way to dereference. The other interesting bits are the NonceManager — which I'll discuss on the next slide — and the explicit status==1 check that prevents reverted transactions from being reported as successful records.""")

# --- Slide P2-10: 3 bugs ---
s = new_slide("Three Engineering Bugs (Real-Chain Frictions)")
add_bullets(s, [
    "Nonce collision under concurrent submission",
    (1, "Symptom: 2 anomalies arriving within 1s → 2nd tx rejected with 'nonce too low'"),
    (1, "Diagnosis: parallel coroutines reading same pending nonce from RPC"),
    (1, "Fix: NonceManager class — asyncio.Lock + locally cached next_nonce + rollback()"),
    "Out-of-gas reverts on logAnomaly",
    (1, "Symptom: every tx reverting with 'out of gas' even though tx hashes existed on chain"),
    (1, "Diagnosis: 180k limit too tight (actual cost ~166k); maxFeePerGas at 30 gwei wasted budget"),
    (1, "Fix: gas → 250k (33% headroom); maxFeePerGas → 5 gwei. Cost dropped 53% (0.000363→0.000167 ETH)"),
    "Failed transactions reported as on-chain records",
    (1, "Symptom: dashboard showed reverted tx as 'logged on-chain'"),
    (1, "Fix: explicit `if rcpt.status != 1: raise` before MSG_BC_LOGGED broadcast"),
], top=Inches(1.35), size=13)
add_footer(s, prs, "Part 2 — slide 10 / 13")
add_speaker_notes(s, """Three bugs that surface only when running against a real public chain. Nonce collision under concurrent submission, out-of-gas because the gas limit was too tight, and reverted-tx-shown-as-success because we trusted the receipt without checking status. Each diagnosed and fixed during the integration session and documented in §6 of the report.""")

# --- Slide P2-11: Rehearsal numbers ---
s = new_slide("Cloud-Services Rehearsal: 2026-04-26 Numbers")
add_table(s, [
    ["Metric", "Result", "Notes"],
    ["Azure IoT Hub messages", "310", "F1 free tier; ~4% of daily 8000 quota"],
    ["S3 anomaly objects", "29", "6 751 bytes total in s3://chenli-cloud-final-2026/anomalies/"],
    ["Sepolia tx confirmed", "18", "Blocks 10 739 301 → 10 739 320, all Status: Success"],
    ["Average gas used", "~166 k", "33% headroom under 250 k limit"],
    ["Cost per tx", "~0.000167 ETH", "Sepolia base fee 1–2 gwei"],
    ["Total wallet spend", "~0.003 ETH", "Free Sepolia ETH; effectively $0 in real money"],
    ["End-to-end click time", "~12 s", "Dominated by single Sepolia block (11–13 s)"],
], left=Inches(0.6), top=Inches(1.6), width=Inches(12.0), height=Inches(4.5), font_size=14)
add_footer(s, prs, "Part 2 — slide 11 / 13")
add_speaker_notes(s, """Yesterday's rehearsal numbers. 310 Azure messages, 29 S3 archives, 18 confirmed Sepolia transactions. Everything paid for in free tier or Sepolia testnet ETH. End-to-end latency from sensor reading to clickable Etherscan link is about 12 seconds, almost entirely dominated by the Sepolia block time.""")

# --- Slide P2-12: Audit chain verification ---
s = new_slide("Audit-Chain Verification: The Killer Property",
               "Pull S3 → recompute SHA-256 → matches on-chain hash exactly")
verify = """1) Etherscan: tx 0xe1ebf113… block 10 739 301 → contract entry recorded
        eventHash = 0x303d3f2164f169f024de0e34ec6438ba5c76da827111e1f3bc8e7cce0bad113a

2) S3 GET object:
   s3://chenli-cloud-final-2026/anomalies/303d3f2164f169...0bad113a.json
   →  {"body_temp":37.37,"device_id":"sim-device-001",
       "hash":"0x303d3f2164f169...0bad113a",
       "heart_rate":86.0,"score":-0.1829,
       "source_iot_node":"iot-sim-device-001",
       "spo2":84.0,"ts":1777249871.0598843}

3) Recompute:  sha256(canonical_json(payload_without_hash))
   = 0x303d3f2164f169f024de0e34ec6438ba5c76da827111e1f3bc8e7cce0bad113a

✅  MATCHES on-chain hash exactly.
   If any byte of the S3 object had been altered, the hashes would diverge.
   This is the concrete, runnable form of the tamper-evident guarantee."""
box = s.shapes.add_textbox(Inches(0.6), Inches(1.95), Inches(12.0), Inches(4.8))
tf = box.text_frame
for i, line in enumerate(verify.split("\n")):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(13)
    if line.strip().startswith("✅"):
        r.font.color.rgb = GREEN
        r.font.bold = True
    elif line.startswith(("1)", "2)", "3)")):
        r.font.color.rgb = ACCENT
        r.font.bold = True
    else:
        r.font.color.rgb = LIGHT_TEXT
add_footer(s, prs, "Part 2 — slide 12 / 13")
add_speaker_notes(s, """This is the slide that proves the architecture actually delivers what the report claims. Take any on-chain entry, fetch the corresponding S3 object, recompute the hash, and verify it matches. I ran this live against transaction 0xe1ebf113 from yesterday's rehearsal — the recomputed hash matches the on-chain hash exactly. If any byte of the S3 record had been altered, the hashes would diverge and the falsification would be detectable. This is the concrete form of the tamper-evident guarantee.""")

# --- Slide P2-13: Midterm response ---
s = new_slide("Addressing Midterm Feedback")
add_bullets(s, [
    "Bootstrap SPOF was the only specific concern raised at midterm",
    "Three layered reasons it's not a SPOF as delivered:",
    (1, "(i) Bootstrap is OFF the data path — sensors/anomalies/tx flow over direct peer connections after join (~200 ms)"),
    (1, "(ii) ANY peer answers MSG_HELLO with MSG_PEERS — not just the role=bootstrap node"),
    (1, "(iii) Multiple bootstraps can run concurrently (similar to Ethereum bootnodes / IPFS)"),
    "Documented in §6.6 of report; demonstrable in live demo by killing the bootstrap process",
    "",
    "Other deltas vs midterm proposal — all upgrades or explicitly disclosed:",
    (1, "Ganache → Sepolia public chain (instructor 4/16 guidance)"),
    (1, "BTP/libp2p → equivalent custom overlay (upstream library readiness)"),
    (1, "S3 expanded to off-chain audit store (was promised; delivered with tamper-evident verification)"),
    (1, "Azure IoT Hub: actively engaged (was promised; delivered)"),
], top=Inches(1.35), size=14)
add_footer(s, prs, "Part 2 — slide 13 / 13")
add_speaker_notes(s, """At midterm the only specific concern raised was: is the bootstrap a single point of failure? The answer, as delivered, is no — and there are three reasons. First, the bootstrap is off the data path; second, any peer can serve discovery; third, multiple bootstraps can coexist. The other midterm-vs-delivered deltas are listed at the bottom — every one is either an upgrade in response to instructor guidance or an explicit disclosure for upstream-library reasons. There's no scope reduction.""")


# ============================================================
# ============ Closing ====================================
# ============================================================

s = new_slide("Conclusion + Q&A")
add_bullets(s, [
    "Part 1: ArchNav legacy app migrated cloud-based to AWS EC2 via Docker — zero source-code changes",
    "Part 2: Decentralized health monitoring system, hybrid Web2 + Web3, 3 active clouds",
    "All five Part-2 design goals (G1-G5) met and demonstrated end-to-end",
    "Real evidence: 310 Azure messages + 29 S3 anomaly archives + 18 confirmed Sepolia transactions + cryptographic audit-chain verification",
    "Engineering integrity: every midterm deviation explicitly disclosed in §6.7 of the final report",
    "",
    "Repository:  github.com/chenli0322/cloud-computing-project (private)",
    "Final report:  part2-health-monitor/docs/final_report.md (~470 lines, 7 sections, 13 references)",
    "",
    "Thank you — Questions?",
], top=Inches(1.35), size=15)
add_footer(s, prs, "End")
add_speaker_notes(s, """Summary: Part 1 delivered cloud-based migration with zero code changes. Part 2 delivered a working decentralized health monitor with three active clouds, real evidence numbers, and a cryptographic audit-chain. Every midterm deviation is disclosed. Open to questions.""")


prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Total slides: {len(prs.slides)}")
